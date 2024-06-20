import pptx
import regex
from fastapi import Depends

from app.services.texts import TextsService
from app.services.documents import DocumentsService
from app.schemas.schemas import *


class PPTXParsingService:

    excluded_texts_re = [
        regex.compile(r'^\d*$'),
        regex.compile(r'^[а-яА-Яa-zA-Z]*$'),
        regex.compile(r'\d{2,4}-\d{2,4}\s?гг')
    ]

    excluded_texts = [
        ' ',
        '',
        '\n'
    ]

    def __init__(
            self,
            texts_service: TextsService = Depends(),
            documents_service: DocumentsService = Depends(),
    ):
        self.current_document = None
        self.current_texts = None
        self.current_document_id = None

        self.texts_service = texts_service
        self.documents_service = documents_service

    async def set_current_docx(self, course_id, type, path_to_file):
        self.current_document = self.load_document(path_to_file)
        doc_res = await self.documents_service.create(
            DocumentPayload(
                courses_id=course_id,
                type=type
            )
        )

        self.current_document_id = doc_res['id']

        self.current_texts = []

        for slide in self.current_document.slides:
            self.current_texts += [slide.notes_slide.notes_text_frame.text]

    @staticmethod
    def load_document(path_to_file):
        return pptx.Presentation(path_to_file)

    def filter_texts(self, texts):
        filtered_texts = []
        for text in texts:
            is_ok = True
            for re in self.excluded_texts_re:
                if re.match(text):
                    is_ok = False
            if text in self.excluded_texts:
                is_ok = False
            if is_ok:
                filtered_texts.append(text.replace('\n', ''))

        return [x for x in filtered_texts if x != '']

    async def get_lecture_text(self):
        self.current_texts = self.filter_texts(self.current_texts)
        for text in self.current_texts:
            await self.texts_service.create(
                TextPayload(
                    document_id=self.current_document_id,
                    competencies_id=None,
                    text=text,
                    type='lecture'
                )
            )
        return len(self.current_texts)
